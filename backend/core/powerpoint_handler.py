"""Handler for PPTX (PowerPoint) files."""

import logging
from pathlib import Path
from typing import Tuple

logger = logging.getLogger(__name__)


class PowerPointHandler:
    """Extract text from PPTX (PowerPoint) files."""

    SUPPORTED_EXTENSIONS = ('.pptx',)

    @staticmethod
    def read_powerpoint(file_path: str) -> Tuple[str, str | None]:
        """
        Extract text from a PPTX file.

        Args:
            file_path: Path to the PPTX file

        Returns:
            (content, error) - content is str, error is None on success
        """
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx not installed. Install with: pip install python-pptx")
            return "", "python-pptx library not installed"

        try:
            prs = Presentation(file_path)
            text_parts = []

            for slide_idx, slide in enumerate(prs.slides, 1):
                text_parts.append(f"--- Slide {slide_idx} ---")

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

                    # Extract table content
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = [cell.text.strip() for cell in row.cells]
                            text_parts.append(" | ".join(row_text))

            content = "\n".join(text_parts)
            logger.info(f"Extracted {len(content)} characters from PowerPoint file: {Path(file_path).name}")
            return content, None

        except FileNotFoundError:
            error_msg = f"File not found: {file_path}"
            logger.error(error_msg)
            return "", error_msg
        except Exception as e:
            error_msg = f"Error reading PowerPoint file: {str(e)}"
            logger.error(error_msg)
            return "", error_msg
